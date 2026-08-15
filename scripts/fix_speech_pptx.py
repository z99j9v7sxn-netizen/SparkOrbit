# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

SRC = Path(r"c:\Users\咸\Desktop\新建 Microsoft PowerPoint 演示文稿 - 副本.pptx")


def get_text_shape(slide):
    for s in slide.shapes:
        if s.has_text_frame:
            return s
    raise RuntimeError("no text shape")


def set_text(slide, text: str, font_size: Optional[float] = None) -> None:
    shape = get_text_shape(slide)
    tf = shape.text_frame
    old_size = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        old_size = tf.paragraphs[0].runs[0].font.size
    size = font_size or (old_size.pt if old_size else 28)

    p0 = tf.paragraphs[0]
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)

    run = p0.add_run()
    run.text = text
    run.font.size = Pt(int(size))
    try:
        p0.alignment = PP_ALIGN.LEFT
    except Exception:
        pass


def main() -> None:
    prs = Presentation(str(SRC))

    t8 = get_text_shape(prs.slides[7]).text
    t8 = t8.replace("踏上我的星轨路径", "踏上我的学习路径")
    set_text(prs.slides[7], t8)
    print("fixed slide 8")

    t14 = get_text_shape(prs.slides[13]).text
    t14 = t14.replace("(伸手示意大屏幕)首先", "(伸手示意大屏幕) 首先")
    set_text(prs.slides[13], t14)
    print("fixed slide 14 spacing")

    for i, slide in enumerate(prs.slides):
        text = get_text_shape(slide).text
        cleaned = text.lstrip("\x0b\x0c\r\n\t ")
        if cleaned != text:
            set_text(slide, cleaned)
            print(f"cleaned leading controls on slide {i + 1}")

    prs.save(str(SRC))
    print("saved")

    prs = Presentation(str(SRC))
    print("total", len(prs.slides))
    vt = "\x0b"
    for i, slide in enumerate(prs.slides, 1):
        text = get_text_shape(slide).text.replace(vt, " ")
        print(f"{i:02d} | {text[:70]}")


if __name__ == "__main__":
    main()
