# -*- coding: utf-8 -*-
"""Align teleprompter PPT slides with design PPT (17 pages)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

SRC = Path(r"c:\Users\咸\Desktop\新建 Microsoft PowerPoint 演示文稿 - 副本.pptx")

TEXTS: dict[int, str | None] = {
    1: (
        "各位老师，好。我是‘爱拼才慧莹’团队的发言人。今天我们带来的作品是《星轨学图 SparkOrbit》——一名学生的认知重生记。"
        "这不是一次常规的答辩，没有冷冰冰的成果汇报；我们更想邀请各位，跟随我的第一视角，推开一扇门——去看看一个“不一样的学习者”，"
        "如何在无数个深夜里，被一串名为 AI 认知孪生的代码，真正“看懂”并温柔陪伴。"
    ),
    2: None,
    3: None,
    4: None,
    5: None,
    6: (
        "“从此，我的知识，变成了我的宇宙。枯燥的学科单元变成了星系，具体的知识点变成了行星。"
        "在艾宾浩斯记忆遗忘规律的运作下，行星开始有了生命：掌握了就会被‘点亮’，大约三天不复习会渐暗，大约七天会触发‘陨石危机’，再久就会暗淡。"
        "星链则把关联知识点簇连成可追逐的成就。那么这些星系是从哪来的？只要老师在后台使用「星系锻造」，一键上传日常用的 PDF 讲义，"
        "就能瞬间为我们开辟出这样一片可探索、可挑战、可复习的知识星图。”"
    ),
    7: (
        "“不是同一份课件发给全班，系统为了我，专门‘打磨’出了六种专属学习资源。"
        "围绕我的画像和薄弱点，六个资源 Agent 深度协同：为我量身定制结构化讲义、思维导图、多层次练习、拓展阅读、"
        "由火山方舟 Seedance 真正生成的教学短视频，以及带注释的代码实操。"
        "更硬核的是，为了防止大模型‘瞎教’，我们首创了资源质量评分系统。所有资源生成后，都必须经过准确性、画像贴合、完整性、幻觉风险四维自动评分，"
        "只有拿了高分才敢呈现在我面前，低分直接在后台自动重试扣留。”"
    ),
    8: (
        "“拿着这些专属资源，我开始踏上我的星轨路径。从诊断起点的黑洞初测出发，到分步任务、打卡推进，再到评估结果回灌后的动态重排——"
        "路径随着我的状态一步步解锁。而我走过的每一步，都凝聚成了这份成长报告：知识掌握率、错题与错因分布、自习区的专注时长，以及六类资源的实际使用情况，"
        "全部从抽象的进步，变成了看得见、摸得着的踏实数据。”"
    ),
    9: (
        "“但是，《星轨学图》认为，我不应该只是一个精密的‘学习机器’。所以除了分数，它也温柔地接住了我的情绪、社交、专注与动力。"
        "在星语树洞里，它默默倾听我的心情日记；在聊天区，班级群聊、私聊和星愿墙让同伴协作变得自然；"
        "而特别是在自习区，我们把‘星座与黄道十二宫’融入其中，让我能选择属于自己的星座房间，在星空的守护下与同伴并肩专注。"
        "这里通过本地摄像头进行视觉督导，但推理全在前端本地运行，视频流绝不上传，只记录分心时长与离开次数这类标量。"
        "休闲区的小游戏、签到商城和桌宠养成，则轻轻托住长期坚持的动力。”"
    ),
    10: (
        "“（深吸一口气）听到这里，各位老师可能觉得，这只是一个孤独学习者的自我沉溺？"
        "现在，让我们切换主观镜头！当我沉浸在自己的星图里时，镜头拉远——原来在我的身后，一直还有一位老师在默默守护！"
        "同一套系统，换一个身份，就是完全不同的战场。在教师端工作台，老师面对的是学情热力图与风险分级、画像复核、AI 教案，"
        "还能用「星系锻造」把 PDF 讲义一键变成可探索课程星系，甚至启动 TimeWarp 沙盘做未来推演。"
        "老师，第一次能在我开口之前，就看懂了我的困惑，将‘因材施教’变成了可控的工作台。”"
    ),
    11: None,
    12: None,
    13: None,
    14: (
        "故事还在升级。为了让“学会了”变得可证明，我们打通了四大增量能力。(伸手示意大屏幕)"
        "首先是四闸掌握策略：行星点亮不再靠答对两三道选择题，而要过「学、练、讲、用」四道关——学过、练过、能讲清、能用上，才算真正 lit。"
        "其次是演武舱：算法过程可以可视化，支持暂停与单步，抽象逻辑摊开看。"
        "再用代码舱：浏览器里直接写 Python、跑测例，服务端结果全绿才记通过。"
        "最后是讯飞数字人与 TTS——需要示范讲解时，虚拟人实时上场，让伴学不只是冷冰冰的文字。"
    ),
    15: (
        "“内容与资产也一起焕新了。一边是校本知识库 Vault：划词与笔记可以沉淀进个人知识资产，还能对接 Obsidian 双链，复习时还能找回——不是一次性聊天。"
        "另一边是 B 站视频深度嵌入：优质课程与科普不用跳出站外，学习页内就能播放、倍速、收藏。"
        "个人积淀与公域视野，终于落在同一条星轨上。”"
    ),
    16: (
        "故事到了尾声。“现在的我，不再是那个‘什么都不会’的学生了。我依然会做错题，但第一次清晰地知道错在哪里；"
        "我依然需要老师，但老师第一次能在无声中拥抱我的困惑。四闸、演武与代码舱、数字人，以及 Vault 与站内视频，让每一次点亮都更可证明、更可追溯。"
        "当然，SparkOrbit 的星图还在不断向外扩张。未来，我们将推进教务系统深度对接、支持遗忘参数的班级个性化配置，甚至打造跨校共享的知识宇宙。"
        "每一颗被点亮的行星，都是科技赐予教育的力量。《星轨学图》，看懂知识，更看懂不一样的你。"
    ),
    17: "我们的汇报完毕，敬请各位评委老师批评指正！”",
}

FONT_SIZES = {
    1: 32,
    6: 28,
    7: 28,
    8: 28,
    9: 26,
    10: 26,
    14: 28,
    15: 28,
    16: 28,
    17: 32,
}


def set_title_text(slide, text: str, font_size: int = 32) -> None:
    shape = None
    for s in slide.shapes:
        if s.has_text_frame:
            shape = s
            break
    if shape is None:
        raise RuntimeError("no text shape on slide")

    tf = shape.text_frame
    try:
        tf.word_wrap = True
    except Exception:
        pass

    p0 = tf.paragraphs[0]
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)

    run = p0.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = False
    try:
        p0.alignment = PP_ALIGN.LEFT
    except Exception:
        pass


def main() -> None:
    prs = Presentation(str(SRC))
    print(f"opened: {SRC.name}, slides={len(prs.slides)}")

    for page in (1, 6, 7, 8, 9, 10):
        text = TEXTS[page]
        if text is not None:
            set_title_text(prs.slides[page - 1], text, FONT_SIZES.get(page, 32))
            print(f"updated slide {page}")

    # Current epilogue / thanks become pages 16 / 17 after reorder
    set_title_text(prs.slides[13], TEXTS[16], FONT_SIZES[16])
    set_title_text(prs.slides[14], TEXTS[17], FONT_SIZES[17])
    print("updated epilogue and thanks (current indices 13/14)")

    layout = prs.slides[0].slide_layout
    s14 = prs.slides.add_slide(layout)
    s15 = prs.slides.add_slide(layout)
    set_title_text(s14, TEXTS[14], FONT_SIZES[14])
    set_title_text(s15, TEXTS[15], FONT_SIZES[15])
    print(f"added two slides at end; total={len(prs.slides)}")

    # Indices after append: 0..12 tech, 13 epilogue, 14 thanks, 15 new14, 16 new15
    # Desired: 0..12, 15, 16, 13, 14
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    desired = list(range(0, 13)) + [15, 16, 13, 14]
    for child in list(sld_id_lst):
        sld_id_lst.remove(child)
    for i in desired:
        sld_id_lst.append(ids[i])

    print(f"reordered; total={len(prs.slides)}")
    prs.save(str(SRC))
    print(f"saved: {SRC}")


if __name__ == "__main__":
    main()
