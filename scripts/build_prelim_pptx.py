# -*- coding: utf-8 -*-
"""Generate SparkOrbit preliminary-round proposal PPTX (12-page product update)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "docs" / "evidence" / "screenshots"
DIAG = ROOT / "docs" / "word" / "img"
OUT = ROOT / "submit" / "SparkOrbit_初赛方案PPT.pptx"
OUT_DESKTOP = Path.home() / "Desktop" / "星轨学图图片版.pptx"

BG_DARK = RGBColor(0x0B, 0x12, 0x20)
BG_LIGHT = RGBColor(0xF4, 0xF7, 0xFB)
TEXT_ON_DARK = RGBColor(0xE8, 0xEE, 0xF7)
TEXT_ON_LIGHT = RGBColor(0x1A, 0x23, 0x32)
ACCENT = RGBColor(0x3D, 0xB8, 0xA0)
WARN = RGBColor(0xE8, 0xA8, 0x38)
MUTED = RGBColor(0x6B, 0x7C, 0x93)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD7, 0xE2)
SOFT = RGBColor(0xE8, 0xF6, 0xF2)
BLUE = RGBColor(0x3A, 0x8F, 0xC0)
TEAL = RGBColor(0x2A, 0x9B, 0x86)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 12

TEAM = "爱拼才慧莹"
SCHOOL = "吉林外国语大学"
DEMO = "https://wikj.online"


def _set_run_font(run, size_pt: float, color: RGBColor, bold: bool = False, name: str = "微软雅黑"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:eastAsian"))
    if ea is None:
        from lxml import etree

        ea = etree.SubElement(rpr, qn("a:eastAsian"))
    ea.set("typeface", name)


def _fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size=18,
    color=TEXT_ON_LIGHT,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size, color, bold=bold)
    return box


def _add_lines(slide, left, top, width, height, lines: list[str], size=16, color=TEXT_ON_LIGHT, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        _set_run_font(run, size, color, bold=(bold_first and i == 0))
    return box


def _notes(slide, text: str):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def _footer(slide, page: int, dark: bool = False):
    color = MUTED if not dark else RGBColor(0x8A, 0x9A, 0xB0)
    _add_textbox(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3), "SparkOrbit 星轨学图", size=10, color=color)
    _add_textbox(
        slide,
        Inches(11.3),
        Inches(7.1),
        Inches(1.5),
        Inches(0.3),
        f"{page} / {TOTAL}",
        size=10,
        color=color,
        align=PP_ALIGN.RIGHT,
    )


def _bg(slide, dark: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill_solid(shape, BG_DARK if dark else BG_LIGHT)
    shape.line.fill.background()


def _accent_bar(slide, left, top, height=Inches(0.35)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    _fill_solid(bar, ACCENT)
    bar.line.fill.background()


def _card(slide, left, top, width, height, fill=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_solid(shape, fill)
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def _title(slide, text: str, dark: bool = False):
    _accent_bar(slide, Inches(0.5), Inches(0.38), Inches(0.42))
    _add_textbox(
        slide,
        Inches(0.7),
        Inches(0.32),
        Inches(12),
        Inches(0.55),
        text,
        size=24,
        color=TEXT_ON_DARK if dark else TEXT_ON_LIGHT,
        bold=True,
    )


def _claim(slide, text: str):
    _add_textbox(slide, Inches(0.7), Inches(0.82), Inches(12), Inches(0.35), text, size=13, color=MUTED)


def _resolve(*candidates: Path) -> Path | None:
    for p in candidates:
        if p and p.exists():
            return p
    return None


def _add_picture(slide, path: Path | None, left, top, width=None, height=None, placeholder_label: str = ""):
    if path and path.exists():
        if width and height:
            slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        elif width:
            slide.shapes.add_picture(str(path), left, top, width=width)
        else:
            slide.shapes.add_picture(str(path), left, top, height=height)
        return True
    ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width or Inches(4), height or Inches(2.5))
    _fill_solid(ph, RGBColor(0xE8, 0xEC, 0xF2))
    ph.line.color.rgb = BORDER
    _add_textbox(
        slide,
        left,
        top + (height or Inches(2.5)) / 2 - Inches(0.2),
        width or Inches(4),
        Inches(0.4),
        placeholder_label or "[截图占位]",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    return False


def _flow_box(slide, left, top, text: str, sub: str):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.85), Inches(1.05))
    _fill_solid(box, WHITE)
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.5)
    _add_textbox(
        slide,
        left + Inches(0.05),
        top + Inches(0.15),
        Inches(1.75),
        Inches(0.4),
        text,
        size=12,
        color=TEXT_ON_LIGHT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        left + Inches(0.05),
        top + Inches(0.55),
        Inches(1.75),
        Inches(0.4),
        sub,
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def _arrow(slide, left, top):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.26), Inches(0.2))
    _fill_solid(arr, ACCENT)
    arr.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    hub = _resolve(SHOTS / "hub_01.png", SHOTS / "resource_media_01.png", DIAG / "diagram-1.png")
    mind = _resolve(SHOTS / "resource_mindmap_01.png", DIAG / "diagram-3.png")
    eval_img = _resolve(SHOTS / "eval_report_01.png")
    path_b = _resolve(SHOTS / "path_before_01.png")
    path_a = _resolve(SHOTS / "path_after_01.png")
    doc = _resolve(SHOTS / "resource_doc_01.png", SHOTS / "tutor_socratic_01.png")
    media = _resolve(SHOTS / "resource_media_01.png")
    hallu = _resolve(SHOTS / "hallu_ticket_teacher_01.png")
    quiz = _resolve(SHOTS / "resource_quiz_01.png")
    tutor = _resolve(SHOTS / "tutor_socratic_01.png")
    diagram = _resolve(DIAG / "diagram-2.png", DIAG / "diagram-4.png", DIAG / "diagram-1.png")
    interview = _resolve(
        SHOTS / "interview_job_text_01.png",
        SHOTS / "interview_academic_01.png",
        SHOTS / "interview_job_followup_01.png",
    )

    # ----- 1 Cover -----
    s = prs.slides.add_slide(blank)
    _bg(s, dark=True)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(1.2), Inches(0.08))
    _fill_solid(accent, ACCENT)
    accent.line.fill.background()
    _add_textbox(s, Inches(0.5), Inches(1.9), Inches(7.6), Inches(0.75), "SparkOrbit 星轨学图", size=38, color=TEXT_ON_DARK, bold=True)
    _add_textbox(
        s,
        Inches(0.5),
        Inches(2.75),
        Inches(7.6),
        Inches(1.15),
        "认知孪生驱动的多智能体学习系统\n作答前预演 · 四闸掌握 · 求职/升学面试闭环",
        size=16,
        color=ACCENT,
    )
    _add_textbox(s, Inches(0.5), Inches(4.05), Inches(7.6), Inches(0.35), "高等教育 · 数据结构 / 机器学习 · Adaptive-LPDS", size=13, color=MUTED)
    _add_lines(
        s,
        Inches(0.5),
        Inches(4.55),
        Inches(7.2),
        Inches(1.6),
        [
            f"参赛团队：{TEAM}",
            f"所属院校：{SCHOOL}",
            f"演示站点：{DEMO}",
            "差异主张：预演仿真 · 四模式 Agent · 四闸掌握 · Shield 终裁",
        ],
        size=13,
        color=RGBColor(0xB0, 0xBC, 0xCE),
    )
    _add_picture(s, hub, Inches(8.2), Inches(1.6), width=Inches(4.5), height=Inches(4.4), placeholder_label="hub_01 / product UI")
    _footer(s, 1, dark=True)
    _notes(
        s,
        "评委老师好，我们是爱拼才慧莹队。星轨学图用多智能体在作答前预演误区，四闸确认学会，并把低置信结果交给教师终裁；出口场景还能走求职/升学面试闭环。",
    )

    # ----- 2 Scene -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "场景来源：高校个性化学习断层")
    _claim(s, "大班统一进度难补救 · Chat 学伴缺闭环 · 生成内容难终裁 → 切计算机/AI 课做可复现 Web 试点")
    left_items = [
        ("大班现实", "同一课件、同一进度，难按人补救"),
        ("AI 助学现状", "多为单聊问答，缺路径规划与评估回灌"),
        ("可信缺口", "生成易幻觉，教师难复核、难终裁"),
    ]
    for i, (h, body) in enumerate(left_items):
        top = Inches(1.3) + Inches(i * 1.15)
        _card(s, Inches(0.5), top, Inches(6.5), Inches(1.05))
        _accent_bar(s, Inches(0.7), top + Inches(0.28), Inches(0.45))
        _add_textbox(s, Inches(1.0), top + Inches(0.18), Inches(5.7), Inches(0.3), h, size=16, color=TEXT_ON_LIGHT, bold=True)
        _add_textbox(s, Inches(1.0), top + Inches(0.55), Inches(5.7), Inches(0.35), body, size=13, color=MUTED)

    _card(s, Inches(7.3), Inches(1.3), Inches(5.5), Inches(3.45))
    _add_textbox(s, Inches(7.55), Inches(1.45), Inches(5.0), Inches(0.35), "落地切入与适用场景", size=15, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(7.55),
        Inches(1.9),
        Inches(5.0),
        Inches(2.65),
        [
            "· 学科：数据结构 / 机器学习（可结构化验证）",
            "· 用户：学生 · 教师 · 教务/管理员",
            "· 大班授课：学情风险筛查 + 人审工单",
            "· 自习助学：画像路径 + 本地督导",
            "· 编程实训：CodeLab + 演武舱",
            "· 校本试点：讲义锻造 → 星系图谱",
            "· 求职/升学：第 7 区模拟面试（第 8 页）",
        ],
        size=12,
        color=TEXT_ON_LIGHT,
    )
    _add_picture(s, mind, Inches(7.3), Inches(4.95), width=Inches(5.5), height=Inches(1.7), placeholder_label="resource_mindmap_01")
    _footer(s, 2)
    _notes(s, "高等教育大班个性化成本高；Chat 学伴停在问答。选计算机类课程便于 Agent、RAG、判题验证。求职与升学作为出口场景，放到第 8 页。")

    # ----- 3 Pain -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "三类用户硬痛点 → 后文用哪一环解")
    _claim(s, "本页只立靶；解法在流程 / Agent / 面试 / 合规页一一对应")
    cols = [
        (
            "学生",
            ["课内：画像粗、资源同质、补救滞后", "出口：面试靠题海，缺多模态评与回流", "→ 六维 + 四闸 + 路径回灌", "→ 第 8 页三官 council + 错题闭环"],
            ACCENT,
        ),
        (
            "教师",
            ["难筛风险与低置信作答", "讲义难结构化、AI 难信", "→ 星系锻造 + Shield 工单", "→ 面试督导与待人审"],
            TEAL,
        ),
        (
            "管理员",
            ["Token / 运维不可见", "缺统一权限与内容台", "→ 用量监测 + 维护模式", "→ /admin/agents 步骤回放"],
            BLUE,
        ),
    ]
    for i, (name, lines, color) in enumerate(cols):
        left = Inches(0.45) + Inches(i * 4.25)
        _card(s, left, Inches(1.25), Inches(4.05), Inches(3.55))
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.25), Inches(4.05), Inches(0.5))
        _fill_solid(head, color)
        head.line.fill.background()
        _add_textbox(s, left, Inches(1.32), Inches(4.05), Inches(0.4), name, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _add_lines(s, left + Inches(0.2), Inches(1.95), Inches(3.65), Inches(2.6), [f"· {x}" for x in lines], size=12, color=TEXT_ON_LIGHT)

    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(5.05), Inches(8.0), Inches(1.55))
    _fill_solid(bar, SOFT)
    bar.line.color.rgb = ACCENT
    _add_textbox(s, Inches(0.7), Inches(5.2), Inches(7.5), Inches(0.35), "映射总览", size=14, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(0.7),
        Inches(5.55),
        Inches(7.5),
        Inches(0.95),
        [
            "学生痛点 → 第 4 页壁垒 · 第 5 页主闭环 · 第 8 页面试",
            "教师痛点 → 第 7 页 Shield 三级 · 第 10 页人机终裁",
            "管理痛点 → 第 6 页管理端 · 第 9 页用量与 Agent 观测",
        ],
        size=12,
        color=TEXT_ON_LIGHT,
    )
    _add_picture(s, eval_img, Inches(8.7), Inches(5.05), width=Inches(4.1), height=Inches(1.55), placeholder_label="eval_report_01")
    _footer(s, 3)
    _notes(s, "逐角色点两句后果；解法不在本页展开，翻页对应。学生出口痛点指向第 8 页面试闭环。")

    # ----- 4 Highlights -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "差异化一览：评委可 10 秒复述的四条壁垒")
    _claim(s, "不是「又一个 Chat 学伴」——预演、掌握协议、人审终裁、四模式 Agent；面试见第 8 页")
    pillars = [
        ("01 预演仿真链", "Teacher→Mirror→Evaluator→PathPlanner", "作答前暴露误区，再进入正式学习", "演示：仿真控制台 SSE"),
        ("02 四闸掌握协议", "学 → 练 → 讲 → 用", "证据齐全才点亮行星，点亮=真学会", "演示：行星门禁进度条"),
        ("03 Shield 人审终裁", "引用 ID + 置信度 + 工单", "低置信/矛盾不终裁，教师可覆盖", "演示：教师待审工单"),
        ("04 四模式 + 七类资源", "workflow / handoff / council / supervisor", "七 Agent 含 Deck 课件；真并行写 AgentStep", "演示：资源工坊 + /admin/agents"),
    ]
    for i, (h, mid, body, demo) in enumerate(pillars):
        left = Inches(0.4) + Inches((i % 2) * 6.45)
        top = Inches(1.25) + Inches((i // 2) * 2.7)
        _card(s, left, top, Inches(6.2), Inches(2.5))
        _add_textbox(s, left + Inches(0.25), top + Inches(0.2), Inches(5.7), Inches(0.4), h, size=16, color=ACCENT, bold=True)
        _add_textbox(s, left + Inches(0.25), top + Inches(0.65), Inches(5.7), Inches(0.35), mid, size=13, color=TEXT_ON_LIGHT, bold=True)
        _add_textbox(s, left + Inches(0.25), top + Inches(1.15), Inches(5.7), Inches(0.55), body, size=13, color=MUTED)
        _add_textbox(s, left + Inches(0.25), top + Inches(1.85), Inches(5.7), Inches(0.35), demo, size=12, color=WARN, bold=True)
    _footer(s, 4)
    _notes(s, "四条壁垒压住全场：预演、四闸、Shield、四模式+七 Agent。模拟面试不在本页抢戏。")

    # ----- 5 Core flow -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "核心流程：采集 → 生成 → 验证 → 反馈")
    _claim(s, "六维可核对 · 七 Agent 可点名 · 四闸是「学会」协议；出口场景走第 8 页面试闭环")
    steps = [
        ("对话画像", "六维抽取"),
        ("资源生成", "七类 Agent"),
        ("路径推送", "画像驱动"),
        ("智能辅导", "苏/费曼"),
        ("效果评估", "雷达回灌"),
        ("教师复核", "待人审"),
    ]
    y = Inches(1.2)
    for i, (t, sub) in enumerate(steps):
        x = Inches(0.35) + Inches(i * 2.15)
        _flow_box(s, x, y, t, sub)
        if i < len(steps) - 1:
            _arrow(s, x + Inches(1.9), y + Inches(0.4))

    _card(s, Inches(0.4), Inches(2.45), Inches(6.3), Inches(1.55), fill=SOFT)
    _add_textbox(s, Inches(0.6), Inches(2.55), Inches(5.9), Inches(0.3), "六维画像（可核对）", size=13, color=ACCENT, bold=True)
    _add_textbox(
        s,
        Inches(0.6),
        Inches(2.95),
        Inches(5.9),
        Inches(0.85),
        "专业背景 · 前置知识 · 认知风格\n易错倾向 · 学习目标 · 时间弹性",
        size=13,
        color=TEXT_ON_LIGHT,
    )
    _card(s, Inches(6.9), Inches(2.45), Inches(5.95), Inches(1.55), fill=SOFT)
    _add_textbox(s, Inches(7.1), Inches(2.55), Inches(5.55), Inches(0.3), "七类资源 Agent（可核对）", size=13, color=ACCENT, bold=True)
    _add_textbox(
        s,
        Inches(7.1),
        Inches(2.95),
        Inches(5.55),
        Inches(0.85),
        "Doc · Mind · Quiz · Read · Media · Deck · Code\nMedia=Seedance 短视频 · Deck=课件/闪卡/PPTX",
        size=13,
        color=TEXT_ON_LIGHT,
    )

    gate = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.15), Inches(12.45), Inches(0.7))
    _fill_solid(gate, BG_DARK)
    gate.line.fill.background()
    _add_textbox(
        s,
        Inches(0.6),
        Inches(4.28),
        Inches(12.1),
        Inches(0.45),
        "点亮门禁（掌握协议）：学 → 练 → 讲 → 用　　证据齐全才算学会　　出口：准备→单轮→总评→回流（第 8 页）",
        size=13,
        color=TEXT_ON_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(s, Inches(0.4), Inches(5.0), Inches(6), Inches(0.3), "路径对比（画像变化前后）", size=11, color=MUTED)
    _add_picture(s, path_b, Inches(0.4), Inches(5.3), width=Inches(6.15), height=Inches(1.5), placeholder_label="path_before_01")
    _add_picture(s, path_a, Inches(6.75), Inches(5.3), width=Inches(6.1), height=Inches(1.5), placeholder_label="path_after_01")
    _footer(s, 5)
    _notes(s, "主链路：对话建档→专属资源→按路径学→辅导不直接给终答→评估改路径；吃不准交给老师。四闸是掌握协议。面试是出口场景。")

    # ----- 6 Product -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "产品形态：三角色 Web · 学生端七区")
    _claim(s, "浏览器可复现 · 教师是治理与终裁 · 主链路在学习区 Dock · 休闲/桌宠答辩 ≤10 秒带过")
    _card(s, Inches(0.45), Inches(1.2), Inches(5.5), Inches(2.55))
    _add_textbox(s, Inches(0.7), Inches(1.32), Inches(5.0), Inches(0.3), "交付形态", size=14, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(0.7),
        Inches(1.7),
        Inches(5.0),
        Inches(1.9),
        [
            "· 形态：浏览器 SPA（非原生 App）",
            "· 接口：REST + SSE + WebSocket",
            "· 部署：Docker Compose 一键",
            "· 可访问：公网 HTTPS 演示站",
        ],
        size=13,
        color=TEXT_ON_LIGHT,
    )
    _card(s, Inches(0.45), Inches(3.9), Inches(5.5), Inches(2.7))
    _add_textbox(s, Inches(0.7), Inches(4.02), Inches(5.0), Inches(0.3), "学生端七大分区", size=14, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(0.7),
        Inches(4.4),
        Inches(5.0),
        Inches(2.05),
        [
            "01 学习区  02 我的星域  03 星语树洞",
            "04 聊天区  05 自习区    06 休闲区",
            "07 模拟面试区（求职 / 升学）",
            "主链路 Dock：星库 · 演武 · CodeLab · 工坊 · 伴学",
        ],
        size=12,
        color=TEXT_ON_LIGHT,
    )
    roles = [
        ("学生端", "七区领航 · 星库划词 · 演武舱 · CodeLab · 面试舱"),
        ("教师端", "星系锻造 · 学情看板 · 人审工单 · TimeWarp · 面试督导"),
        ("管理端", "用户 / 内容 / Token 用量 / 维护模式 / Agent 回放"),
    ]
    for i, (h, b) in enumerate(roles):
        top = Inches(1.2) + Inches(i * 1.15)
        _card(s, Inches(6.2), top, Inches(6.6), Inches(1.05))
        _add_textbox(s, Inches(6.45), top + Inches(0.12), Inches(6.15), Inches(0.3), h, size=14, color=ACCENT, bold=True)
        _add_textbox(s, Inches(6.45), top + Inches(0.5), Inches(6.15), Inches(0.4), b, size=13, color=TEXT_ON_LIGHT)
    _add_textbox(s, Inches(6.2), Inches(4.7), Inches(6.6), Inches(0.3), "校本 B2B2C（讲义→星系）+ 可选 B2C / Token", size=12, color=MUTED, bold=True)
    _add_picture(s, doc, Inches(6.2), Inches(5.05), width=Inches(3.15), height=Inches(1.55), placeholder_label="resource_doc_01")
    _add_picture(s, tutor, Inches(9.5), Inches(5.05), width=Inches(3.3), height=Inches(1.55), placeholder_label="tutor_socratic_01")
    _footer(s, 6)
    _notes(s, "打开链接即可验；三角色注册分流。主分看画像—资源—路径—辅导—评估；休闲功能十秒带过。第七区是模拟面试。")

    # ----- 7 Agents -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "Agent 能力：四种真编排，不是单模型包办")
    _claim(s, "LangGraph 只用于 handoff；资源/准备走 workflow DAG；勿称全系统都是 LangGraph")

    modes = [
        ("workflow", "资源工坊 / 面试准备", "C2 三组 DAG 真并行，写 AgentStep"),
        ("handoff", "镜像预演 / 面试单轮", "LangGraph 真正 astream"),
        ("council", "平行宇宙 / 面试总评", "多策略或三官并行后汇总"),
        ("supervisor", "伴学辅导", "意图分类 → 路径 / 资源 / 闪卡"),
    ]
    for i, (mode, scene, act) in enumerate(modes):
        left = Inches(0.4) + Inches((i % 4) * 3.2)
        _card(s, left, Inches(1.2), Inches(3.05), Inches(1.7))
        _add_textbox(s, left + Inches(0.12), Inches(1.3), Inches(2.8), Inches(0.3), mode, size=14, color=ACCENT, bold=True)
        _add_textbox(s, left + Inches(0.12), Inches(1.65), Inches(2.8), Inches(0.45), scene, size=11, color=TEXT_ON_LIGHT, bold=True)
        _add_textbox(s, left + Inches(0.12), Inches(2.15), Inches(2.8), Inches(0.55), act, size=11, color=MUTED)

    _card(s, Inches(0.4), Inches(3.1), Inches(8.2), Inches(1.15))
    _add_textbox(s, Inches(0.6), Inches(3.18), Inches(7.8), Inches(0.28), "代表链 · 仿真 handoff", size=13, color=ACCENT, bold=True)
    _add_textbox(
        s,
        Inches(0.6),
        Inches(3.5),
        Inches(7.8),
        Inches(0.6),
        "Teacher → Mirror → Evaluator → PathPlanner　　作答前预演误区，再定制路径（SSE 可观察）",
        size=12,
        color=TEXT_ON_LIGHT,
    )
    _card(s, Inches(0.4), Inches(4.4), Inches(8.2), Inches(1.15))
    _add_textbox(s, Inches(0.6), Inches(4.48), Inches(7.8), Inches(0.28), "代表链 · 资源 workflow", size=13, color=ACCENT, bold=True)
    _add_textbox(
        s,
        Inches(0.6),
        Inches(4.8),
        Inches(7.8),
        Inches(0.6),
        "Doc / Mind / Quiz / Read / Media / Deck / Code　　同组独立 DB session 并行，附质量评分",
        size=12,
        color=TEXT_ON_LIGHT,
    )
    _card(s, Inches(0.4), Inches(5.7), Inches(8.2), Inches(1.0))
    _add_textbox(
        s,
        Inches(0.6),
        Inches(5.82),
        Inches(7.8),
        Inches(0.75),
        "Shield 三级：ID+RAG 依据 → Evaluator 独立引用/置信 → 低置信进工单　　苏/费曼辅导 · AI 不终裁\n观测：/admin/agents 按 user_id 回放 AgentStep",
        size=12,
        color=TEXT_ON_LIGHT,
    )
    _add_picture(s, media, Inches(8.85), Inches(3.1), width=Inches(3.95), height=Inches(1.7), placeholder_label="resource_media_01")
    _add_picture(s, hallu, Inches(8.85), Inches(4.95), width=Inches(3.95), height=Inches(1.75), placeholder_label="hallu_ticket_teacher_01")
    _footer(s, 7)
    _notes(s, "四种 mode 对上 AGENTS.md。仿真用 LangGraph astream；资源用 DAG 真并行。禁止说全系统都是 LangGraph。")

    # ----- 8 Interview -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "模拟面试：课内闭环的出口场景")
    _claim(s, "求职 / 升学共用四模式 Agent；评完回流错题与复习卡，不是另一套聊天框")

    _card(s, Inches(0.4), Inches(1.2), Inches(4.05), Inches(2.35))
    _add_textbox(s, Inches(0.6), Inches(1.32), Inches(3.7), Inches(0.3), "双场景", size=14, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(0.6),
        Inches(1.75),
        Inches(3.7),
        Inches(1.6),
        [
            "· 求职校招：技术 / HR / 业务三官",
            "· 升学复试：学科 / 素质 / 科研三官",
            "· 入口：学生端第 7 区面试舱",
            "· 教师端可下发任务并督导",
        ],
        size=12,
        color=TEXT_ON_LIGHT,
    )
    stages = [
        ("1 准备 workflow", "JobAnalyst ∥ ProfileParser\n→ QuestionPlanner → Q-* 并行"),
        ("2 单轮 handoff", "AnswerAggregator →\nMultimodalScorer → FollowUp"),
        ("3 总评 council", "三官 asyncio.gather\n→ CouncilSummarizer"),
    ]
    for i, (h, b) in enumerate(stages):
        left = Inches(4.65) + Inches(i * 2.85)
        _card(s, left, Inches(1.2), Inches(2.7), Inches(2.35))
        _add_textbox(s, left + Inches(0.12), Inches(1.32), Inches(2.45), Inches(0.55), h, size=13, color=ACCENT, bold=True)
        _add_textbox(s, left + Inches(0.12), Inches(1.95), Inches(2.45), Inches(1.35), b, size=12, color=TEXT_ON_LIGHT)

    _card(s, Inches(0.4), Inches(3.7), Inches(8.2), Inches(2.9))
    _add_textbox(s, Inches(0.6), Inches(3.82), Inches(7.8), Inches(0.3), "产品能力与回流", size=14, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(0.6),
        Inches(4.25),
        Inches(7.8),
        Inches(2.15),
        [
            "· 数字人播报 + 讯飞 IAT 听写；弱项自动插入一层追问",
            "· 五维雷达报告；低分回流错题本 / 复习卡 / 复盘资源包",
            "· 管理端 scene=interview 可按步骤回放",
            "· 降级：关麦 / 关视觉仍可用文本完成，不假装全模态必开",
        ],
        size=13,
        color=TEXT_ON_LIGHT,
    )
    _add_picture(
        s,
        interview,
        Inches(8.85),
        Inches(3.7),
        width=Inches(3.95),
        height=Inches(2.9),
        placeholder_label="现场点验：第7区面试舱",
    )
    _footer(s, 8)
    _notes(s, "面试不是另一套聊天框，而是把四模式 Agent 接到就业/升学出口，低分回流学习闭环。")

    # ----- 9 Stack -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "工具 · 数据 · 模型：按任务路由")
    _claim(s, "选型服从任务；校本 RAG 管忠实度；敏感视觉能本地则本地")
    blocks = [
        (
            "工具",
            [
                "前端：Vue3 · Vite · Three.js · TF.js",
                "后端：FastAPI · SQLAlchemy · SSE/WS",
                "ChromaDB RAG · LangGraph（仅 handoff）",
                "AgentRun / AgentStep 落库",
                "管理端 /admin/agents 步骤回放",
            ],
        ),
        (
            "数据",
            [
                "六维画像与学习事件",
                "掌握度 / 四闸证据 / 路径评估",
                "校本教材向量（sources 可展示）",
                "面试会话与媒体（学生可删）",
                "Vault 笔记 · 用量日志 · 密码哈希",
            ],
        ),
        (
            "模型（任务路由）",
            [
                "DeepSeek：推理、评分、归因",
                "讯飞：中文 / IAT / ISE / 数字人",
                "火山 Seedance：教学短视频",
                "通义：分身图生图；视觉理解仪态",
                "面试模态可降级为纯文本完成",
            ],
        ),
    ]
    for i, (h, lines) in enumerate(blocks):
        left = Inches(0.4) + Inches(i * 4.25)
        _card(s, left, Inches(1.2), Inches(4.05), Inches(3.7))
        _add_textbox(s, left + Inches(0.2), Inches(1.35), Inches(3.65), Inches(0.35), h, size=15, color=ACCENT, bold=True)
        _add_lines(s, left + Inches(0.2), Inches(1.85), Inches(3.65), Inches(2.8), [f"· {x}" for x in lines], size=12, color=TEXT_ON_LIGHT)

    foot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.15), Inches(8.3), Inches(1.4))
    _fill_solid(foot, SOFT)
    foot.line.color.rgb = ACCENT
    _add_textbox(
        s,
        Inches(0.65),
        Inches(5.4),
        Inches(7.9),
        Inches(0.95),
        "数据流原则：外发仅 prompt / 必要音视频；\n督导摄像头视频流不出浏览器（仅分心/离开等标量落库）",
        size=14,
        color=TEXT_ON_LIGHT,
        bold=True,
    )
    _add_picture(s, diagram, Inches(8.95), Inches(5.15), width=Inches(3.85), height=Inches(1.4), placeholder_label="diagram-2")
    _footer(s, 9)
    _notes(s, "单一模型扛不住语音、视频和长文本教学。按任务选型；校本 RAG 管忠实度。LangGraph 仅 handoff。")

    # ----- 10 Compliance -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "合规边界：能做什么，不能假装做什么")
    _claim(s, "主动划界比吹「全面合规」加分；算法建议可误导 → 人机协同终裁")
    _card(s, Inches(0.4), Inches(1.2), Inches(6.15), Inches(4.35))
    head_l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.2), Inches(6.15), Inches(0.5))
    _fill_solid(head_l, ACCENT)
    head_l.line.fill.background()
    _add_textbox(s, Inches(0.4), Inches(1.28), Inches(6.15), Inches(0.35), "已实现（可演示）", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_lines(
        s,
        Inches(0.65),
        Inches(1.9),
        Inches(5.7),
        Inches(3.4),
        [
            "· 自习督导 TF.js 本地推理，视频不出浏览器",
            "· RBAC 三角色鉴权；密钥 .env 不进仓库",
            "· Shield 风控；低置信不终裁，教师可覆盖",
            "· 学生删除面试会话 → 媒体与报告一并消失",
            "· 定位：竞赛演示 + 校本试点基线",
        ],
        size=13,
        color=TEXT_ON_LIGHT,
    )
    _card(s, Inches(6.8), Inches(1.2), Inches(6.05), Inches(4.35))
    head_r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.2), Inches(6.05), Inches(0.5))
    _fill_solid(head_r, WARN)
    head_r.line.fill.background()
    _add_textbox(s, Inches(6.8), Inches(1.28), Inches(6.05), Inches(0.35), "未宣称 / 后续增强", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_lines(
        s,
        Inches(7.05),
        Inches(1.9),
        Inches(5.55),
        Inches(2.2),
        [
            "· 完整不可篡改审计库、等保级合规",
            "· 未成年人监护人明示同意（生产增强）",
            "· PDF 页码级强制引用（规划中）",
            "· LMS 已对接（规划中，不写已上线）",
        ],
        size=13,
        color=TEXT_ON_LIGHT,
    )
    _add_picture(s, hallu, Inches(7.15), Inches(3.7), width=Inches(5.4), height=Inches(1.6), placeholder_label="hallu_ticket_teacher_01")
    principle = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.75), Inches(12.45), Inches(0.8))
    _fill_solid(principle, BG_DARK)
    principle.line.fill.background()
    _add_textbox(
        s,
        Inches(0.6),
        Inches(5.95),
        Inches(12.1),
        Inches(0.45),
        "原则：算法建议可误导 → 人机协同终裁；隐私能本地则本地；面试媒体可删",
        size=14,
        color=TEXT_ON_DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _footer(s, 10)
    _notes(s, "主动说清边界。点本地督导、教师终裁、面试删除三处硬证据。不宣称等保、LMS、监护人同意。")

    # ----- 11 Roadmap -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "后续落地：先嵌教学流程，再谈跨校生态")
    _claim(s, "近期做实校本；中期对接 LMS；长期本地化呼应合规页；面试先校本就业/升学试点")
    phases = [
        ("近期 · 6 个月内", ACCENT, ["扩充校本 RAG，收紧行星引用边界", "遗忘/复习阈值班级可配，自动派发", "移动端降级（星图 2D、督导轻量）"]),
        ("中期 · 6–18 个月", BLUE, ["对接校内 LMS / 花名册与成绩回写", "「预演 vs 真实作答」准确率校准", "Agent 决策可解释报告"]),
        ("长期 · 18 个月以上", WARN, ["跨校星系共创（隐私合规前提）", "口语/实验/编程多模态闭环加深", "校内本地化或隐私计算，数据不出校"]),
    ]
    for i, (title, color, lines) in enumerate(phases):
        top = Inches(1.2) + Inches(i * 1.55)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), top, Inches(0.15), Inches(1.4))
        _fill_solid(bar, color)
        bar.line.fill.background()
        _card(s, Inches(0.7), top, Inches(12.1), Inches(1.4))
        _add_textbox(s, Inches(1.0), top + Inches(0.15), Inches(11.5), Inches(0.32), title, size=16, color=color, bold=True)
        _add_textbox(s, Inches(1.0), top + Inches(0.58), Inches(11.5), Inches(0.6), "　|　".join(lines), size=13, color=TEXT_ON_LIGHT)
    foot_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(5.9), Inches(12.4), Inches(0.75))
    _fill_solid(foot_r, SOFT)
    foot_r.line.color.rgb = ACCENT
    _add_textbox(
        s,
        Inches(0.7),
        Inches(6.08),
        Inches(12.0),
        Inches(0.45),
        "面试闭环：先做校本就业 / 升学辅导试点，再谈跨校题库。",
        size=14,
        color=TEXT_ON_LIGHT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _footer(s, 11)
    _notes(s, "落地优先嵌进教学流程，再谈生态；长期本地化呼应合规页。面试先校本试点。")

    # ----- 12 Progress -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _title(s, "进度与复现：按主链路点验")
    _claim(s, "画像 → 资源 → 路径 → 辅导 → 评估 → 教师工单；出口再点第 7 区面试（勿先点休闲区）")
    _card(s, Inches(0.4), Inches(1.2), Inches(6.15), Inches(3.35))
    _add_textbox(s, Inches(0.65), Inches(1.35), Inches(5.7), Inches(0.35), "已完成", size=15, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(0.65),
        Inches(1.8),
        Inches(5.7),
        Inches(2.55),
        [
            "· 学生端七区 + 三角色 Web 可运行",
            "· 七类 Agent · 四模式可回放 AgentStep",
            "· 模拟面试主路径（求职 / 升学）",
            "· 四闸门禁 · 星库 · 演武舱 · CodeLab",
            "· Shield 人审 · Docker Compose + 公网站",
        ],
        size=13,
        color=TEXT_ON_LIGHT,
    )
    _card(s, Inches(6.8), Inches(1.2), Inches(6.05), Inches(3.35))
    _add_textbox(s, Inches(7.05), Inches(1.35), Inches(5.55), Inches(0.35), "验证入口", size=15, color=ACCENT, bold=True)
    _add_lines(
        s,
        Inches(7.05),
        Inches(1.8),
        Inches(5.55),
        Inches(2.55),
        [
            f"· 地址：{DEMO}",
            "· 角色：student001 / teacher001 / admin001",
            "· 密码：见部署说明（本页不印明文）",
            "· 证据：docs/evidence/ 案例与截图",
            "· 观测：/admin/agents · scene=interview",
        ],
        size=13,
        color=TEXT_ON_LIGHT,
    )
    _add_picture(s, eval_img or quiz, Inches(0.4), Inches(4.75), width=Inches(4.0), height=Inches(1.9), placeholder_label="eval_report_01")
    _add_picture(s, media, Inches(4.6), Inches(4.75), width=Inches(4.0), height=Inches(1.9), placeholder_label="resource_media_01")
    thanks = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.85), Inches(4.75), Inches(4.0), Inches(1.9))
    _fill_solid(thanks, BG_DARK)
    thanks.line.fill.background()
    _add_textbox(s, Inches(8.85), Inches(5.05), Inches(4.0), Inches(0.5), "谢谢评委老师", size=20, color=TEXT_ON_DARK, bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(
        s,
        Inches(8.95),
        Inches(5.6),
        Inches(3.8),
        Inches(0.8),
        "预演 · 四闸 · Shield · 面试闭环\n欢迎按主链路点验",
        size=12,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    _footer(s, 12)
    _notes(s, "请按主链路点击验证；被问休闲功能 10 秒带过，拉回闭环。出口场景再点第 7 区面试。")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    try:
        prs.save(str(OUT_DESKTOP))
        desktop_msg = str(OUT_DESKTOP)
    except Exception as exc:
        desktop_msg = f"(desktop save skipped: {exc})"
    print(f"Wrote {OUT}")
    print(f"Also wrote {desktop_msg}")
    print(
        "Images:",
        f"hub={hub.name if hub else None}",
        f"media={media.name if media else None}",
        f"hallu={hallu.name if hallu else None}",
        f"path={bool(path_b and path_a)}",
        f"eval={eval_img.name if eval_img else None}",
        f"interview={interview.name if interview else None}",
    )


if __name__ == "__main__":
    build()
