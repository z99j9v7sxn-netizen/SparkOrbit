# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

SRC = Path(r"c:\Users\咸\Desktop\新建 Microsoft PowerPoint 演示文稿 - 副本.pptx")


def get_text_shape(slide):
    for s in slide.shapes:
        if s.has_text_frame:
            return s
    raise RuntimeError("no text shape")


def set_clean_text(slide, text: str, font_size: Optional[float] = None) -> None:
    shape = get_text_shape(slide)
    tf = shape.text_frame
    old_size = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        old_size = tf.paragraphs[0].runs[0].font.size
    size = font_size or (old_size.pt if old_size else 28)

    # Keep only first paragraph; strip all runs/breaks inside it
    p0 = tf.paragraphs[0]
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

    p_elem = p0._p
    for child in list(p_elem):
        tag = child.tag
        if tag.endswith("}r") or tag.endswith("}br") or tag.endswith("}fld"):
            p_elem.remove(child)

    run = p0.add_run()
    run.text = text.lstrip("\x0b\x0c\r\n\t ")
    run.font.size = Pt(int(size))
    try:
        p0.alignment = PP_ALIGN.LEFT
    except Exception:
        pass


def main() -> None:
    prs = Presentation(str(SRC))
    for i, slide in enumerate(prs.slides):
        shape = get_text_shape(slide)
        raw = shape.text
        cleaned = raw.replace("\x0b", "").replace("\x0c", "").strip()
        # collapse accidental double spaces from removed breaks
        while "  " in cleaned:
            cleaned = cleaned.replace("  ", " ")
        set_clean_text(slide, cleaned)
        after = get_text_shape(slide).text
        print(f"{i + 1:02d} start_repr={after[:12]!r}")

    prs.save(str(SRC))
    print("saved", SRC)


if __name__ == "__main__":
    main()
