# -*- coding: utf-8 -*-
"""生成 GOAI 演示提词 PPT：横版 16:9，一页一段大字，白底深字。

输入：docs/evidence/录屏演讲稿_GOAI提词版.md（## 第 N 页 开头的 18 页正文）
输出：桌面 SparkOrbit_GOAI演示提词.pptx（不覆盖旧文件）
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

SRC = Path(r"c:\Users\咸\Desktop\project\docs\evidence\录屏演讲稿_GOAI提词版.md")
OUT = Path(r"c:\Users\咸\Desktop\SparkOrbit_GOAI演示提词.pptx")

FONT = "微软雅黑"
TEXT_COLOR = RGBColor(0x1A, 0x23, 0x32)


def parse_pages(text: str) -> list[str]:
    """按 '## 第 N 页' 切分，返回每页正文（标题后到下一个 '---' 前的段落文本列表）。"""
    pages: list[str] = []
    blocks = re.split(r"^## 第 \d+ 页.*$", text, flags=re.M)
    for block in blocks[1:]:
        # 去掉块内第一个 '---' 之后的内容（那是分隔符/时间轴）
        body = block.split("\n---\n", 1)[0]
        lines = [ln.rstrip() for ln in body.split("\n")]
        # 去掉开头/结尾空行，保留内部段落
        while lines and not lines[0].strip():
            lines.pop(0)
        while lines and not lines[-1].strip():
            lines.pop()
        pages.append("\n".join(lines))
    return pages


def pick_size(length: int) -> int:
    if length <= 90:
        return 36
    if length <= 170:
        return 32
    if length <= 230:
        return 28
    return 24


def build(pages: list[str]) -> None:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    blank = prs.slide_layouts[6]  # blank layout

    for idx, body in enumerate(pages, 1):
        slide = prs.slides.add_slide(blank)

        tb = slide.shapes.add_textbox(
            Emu(914400), Emu(548640), Emu(10287000), Emu(5760720)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        paragraphs = body.split("\n")
        # 去掉首尾空段
        while paragraphs and not paragraphs[0].strip():
            paragraphs.pop(0)
        while paragraphs and not paragraphs[-1].strip():
            paragraphs.pop()

        total_len = sum(len(p) for p in paragraphs)
        size = pick_size(total_len)

        for i, para_text in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)
            run = p.add_run()
            run.text = para_text
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = TEXT_COLOR

    prs.save(str(OUT))
    print("saved", OUT)
    print("pages", len(pages))


def main() -> None:
    text = SRC.read_text(encoding="utf-8")
    pages = parse_pages(text)
    if len(pages) != 18:
        raise SystemExit(f"预期 18 页，实际解析出 {len(pages)} 页")
    for i, p in enumerate(pages, 1):
        print(f"{i:02d} | {len(p)} 字 | {p[:28]!r}")
    build(pages)


if __name__ == "__main__":
    main()
