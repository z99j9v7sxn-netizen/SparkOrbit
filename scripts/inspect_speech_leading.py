# -*- coding: utf-8 -*-
from pathlib import Path
from pptx import Presentation

SRC = Path(r"c:\Users\咸\Desktop\新建 Microsoft PowerPoint 演示文稿 - 副本.pptx")
prs = Presentation(str(SRC))

for idx in (0, 9):
    shape = next(s for s in prs.slides[idx].shapes if s.has_text_frame)
    text = shape.text
    print(f"=== slide {idx + 1} ===")
    print("repr start:", repr(text[:20]))
    print("ords:", [ord(c) for c in text[:8]])
    for pi, p in enumerate(shape.text_frame.paragraphs):
        for ri, r in enumerate(p.runs):
            print(f"  p{pi}r{ri} repr={r.text[:20]!r} ords={[ord(c) for c in r.text[:5]]}")
