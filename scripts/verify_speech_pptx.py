# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation

speech = Path(r"c:\Users\咸\Desktop\新建 Microsoft PowerPoint 演示文稿 - 副本.pptx")
design = Path(r"c:\Users\咸\Desktop\SparkOrbit_星轨学图_优化美化版222(5).pptx")

chapters = [
    "1 开场/认知重生记",
    "2 PROLOGUE 五个困境",
    "3 领航台六区",
    "4 CHAPTER1 画像",
    "5 CHAPTER2 镜像预演",
    "6 CHAPTER3 星系/衰减",
    "7 CHAPTER4 六种资源",
    "8 CHAPTER5 路径/成长",
    "9 CHAPTER6 情绪专注社交",
    "10 教师端",
    "11 管理端",
    "12 三端总览",
    "13 技术底座/防幻觉",
    "14 NEW FEATURES 四闸等",
    "15 FEATURE UPDATE Vault/B站",
    "16 EPILOGUE",
    "17 感谢",
]

keywords = [
    ["认知重生", "爱拼才慧莹", "第一视角"],
    ["刷了很多题", "越错越沉默", "只有分数"],
    ["星轨领航台", "六个"],
    ["数字画像", "六个维度"],
    ["另一个我", "Teacher", "Mirror"],
    ["星系", "行星", "陨石", "星系锻造"],
    ["六种", "Seedance", "四维"],
    ["学习路径", "成长报告", "动态重排"],
    ["星语树洞", "自习区", "聊天区", "休闲区"],
    ["教师", "TimeWarp", "星系锻造"],
    ["管理端", "Token", "RBAC"],
    ["三重身份", "一套系统"],
    ["多模型", "防幻觉", "待审"],
    ["四闸", "演武舱", "代码舱", "数字人"],
    ["Vault", "Obsidian", "B 站"],
    ["现在的我", "什么都不会", "跨校"],
    ["批评指正", "汇报完毕"],
]

prs = Presentation(str(speech))
dprs = Presentation(str(design))
print(f"speech={len(prs.slides)} design={len(dprs.slides)}")
assert len(prs.slides) == 17 and len(dprs.slides) == 17

all_ok = True
for i, slide in enumerate(prs.slides, 1):
    text = "".join(s.text for s in slide.shapes if s.has_text_frame)
    missing = [k for k in keywords[i - 1] if k not in text]
    status = "OK" if not missing else f"MISSING {missing}"
    if missing:
        all_ok = False
    print(f"{i:02d} {chapters[i - 1]} | {status}")
    print(f"   {text[:100]}")

print("RESULT:", "ALL_OK" if all_ok else "HAS_GAPS")
