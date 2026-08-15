"""将分镜 slides 导出为可下载的教学 PPTX。"""
from __future__ import annotations

import logging
import uuid
from pathlib import Path
from typing import Any

from app.services.deck_themes import resolve_theme

logger = logging.getLogger(__name__)

_MEDIA_GENERATED = Path(__file__).resolve().parents[1] / "static" / "media" / "generated"


def _rgb(rgb: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def _fill_slide(slide, rgb: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(rgb)


def _add_bar(slide, theme: dict[str, Any], *, top: float, height: float) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), Inches(13.333), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme["bar"])
    shape.line.fill.background()


def _set_run_color(paragraph, rgb: tuple[int, int, int]) -> None:
    for run in paragraph.runs:
        run.font.color.rgb = _rgb(rgb)


def export_deck_pptx(
    *,
    title: str,
    slides: list[dict[str, Any]],
    planet_slug: str = "deck",
    theme_id: str = "orbit",
) -> str:
    """生成 pptx，返回 /static/media/generated/... 路径。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("未安装 python-pptx，无法导出课件") from exc

    theme = resolve_theme(theme_id)
    _MEDIA_GENERATED.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    _fill_slide(cover, theme["bg"])
    _add_bar(cover, theme, top=0, height=0.18)
    _add_bar(cover, theme, top=7.32, height=0.18)
    box = cover.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title or "教学课件"
    p.font.size = Pt(36)
    p.font.bold = True
    _set_run_color(p, theme["title"])
    sub = tf.add_paragraph()
    sub.text = f"星轨学图 · {planet_slug} · {theme['name']}"
    sub.font.size = Pt(16)
    _set_run_color(sub, theme["body"])

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        s = prs.slides.add_slide(blank)
        _fill_slide(s, theme["bg"])
        _add_bar(s, theme, top=0, height=0.12)

        title_box = s.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(1))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = str(slide.get("title") or "要点")
        tp.font.size = Pt(28)
        tp.font.bold = True
        _set_run_color(tp, theme["title"])

        bullets = slide.get("bullet_points") or []
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        body = s.shapes.add_textbox(Inches(0.9), Inches(1.45), Inches(11.5), Inches(5.2))
        bf = body.text_frame
        bf.word_wrap = True
        bf.clear()
        first = True
        lines = bullets[:8] or [str(slide.get("narration") or "")[:200] or "要点"]
        size = Pt(20) if bullets else Pt(18)
        for b in lines:
            para = bf.paragraphs[0] if first else bf.add_paragraph()
            first = False
            para.text = str(b)
            para.level = 0
            para.font.size = size
            _set_run_color(para, theme["body"])

        notes = s.notes_slide.notes_text_frame
        notes.text = str(slide.get("narration") or "")

    filename = f"deck_{planet_slug}_{uuid.uuid4().hex[:10]}.pptx"
    dest = _MEDIA_GENERATED / filename
    prs.save(str(dest))
    logger.info("exported deck pptx %s slides=%s theme=%s", dest.name, len(slides), theme["id"])
    return f"/static/media/generated/{filename}"
